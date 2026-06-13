import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation and Widescreen Mode (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Constants
COLOR_BG = RGBColor(228, 228, 230)
COLOR_CARD = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(212, 212, 216)
COLOR_TEXT_PRIMARY = RGBColor(24, 24, 27)
COLOR_TEXT_SECONDARY = RGBColor(63, 63, 70)
COLOR_TEXT_MUTED = RGBColor(113, 113, 122)
COLOR_ACCENT_BLUE = RGBColor(29, 78, 216)
COLOR_ACCENT_GREEN = RGBColor(4, 120, 87)
COLOR_ACCENT_PINK = RGBColor(190, 24, 93)

# Blank layout helper
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def format_run(run, font_name="Segoe UI", font_size=Pt(14), color=COLOR_TEXT_SECONDARY, bold=False, italic=False):
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic

def add_slide_header(slide, category, title):
    # Category (small uppercase text)
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = Inches(0)
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    format_run(p_cat.runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_ACCENT_BLUE, bold=True)
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = Inches(0)
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    format_run(p_title.runs[0], font_name="Segoe UI", font_size=Pt(28), color=COLOR_TEXT_PRIMARY, bold=True)

def add_example_box(slide, left, top, width, height, example_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(244, 244, 245)
    shape.line.color.rgb = COLOR_ACCENT_BLUE
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    
    p0 = tf.paragraphs[0]
    p0.text = "REAL-WORLD EXAMPLE"
    p0.space_after = Pt(3)
    format_run(p0.runs[0], font_name="Segoe UI", font_size=Pt(9.5), color=COLOR_ACCENT_BLUE, bold=True)
    
    p1 = tf.add_paragraph()
    p1.text = example_text
    format_run(p1.runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_SECONDARY, italic=True)

def add_card_box(slide, left, top, width, height, card_title, border_color=COLOR_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    p0 = tf.paragraphs[0]
    p0.text = card_title
    p0.space_after = Pt(12)
    format_run(p0.runs[0], font_name="Segoe UI", font_size=Pt(16), color=COLOR_TEXT_PRIMARY, bold=True)
    
    return tf

def add_bullets(slide, left, top, width, height, bullets):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    
    for idx, (title, desc) in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        
        run_arrow = p.add_run()
        run_arrow.text = "→  "
        format_run(run_arrow, font_name="Segoe UI", font_size=Pt(13), color=COLOR_ACCENT_BLUE, bold=True)
        
        run_title = p.add_run()
        run_title.text = title + ": "
        format_run(run_title, font_name="Segoe UI", font_size=Pt(12.5), color=COLOR_TEXT_PRIMARY, bold=True)
        
        run_desc = p.add_run()
        run_desc.text = desc
        format_run(run_desc, font_name="Segoe UI", font_size=Pt(12.5), color=COLOR_TEXT_SECONDARY)

def draw_arrow_line(slide, x1, y1, x2, y2, color):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml import parse_xml
    
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    
    # Add arrowhead to tailEnd (end of path)
    ln = connector.line._get_or_add_ln()
    tailEnd = parse_xml(
        '<a:tailEnd type="arrow" w="med" len="med" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    )
    ln.append(tailEnd)
    return connector

def parse_code_line(p, line, keywords):
    # Simple inline parser for python formatting
    char_idx = 0
    words = []
    current = ""
    in_string = False
    string_char = None
    
    while char_idx < len(line):
        c = line[char_idx]
        
        if c == '#' and not in_string:
            if current:
                words.append((current, 'normal'))
                current = ""
            words.append((line[char_idx:], 'comment'))
            break
            
        if c in ('"', "'") and (char_idx == 0 or line[char_idx-1] != '\\'):
            if in_string:
                if c == string_char:
                    current += c
                    words.append((current, 'string'))
                    current = ""
                    in_string = False
                    string_char = None
                else:
                    current += c
            else:
                if current:
                    words.append((current, 'normal'))
                    current = ""
                in_string = True
                string_char = c
                current += c
            char_idx += 1
            continue
            
        if in_string:
            current += c
        else:
            if c in (' ', '(', ')', '[', ']', '{', '}', '=', ',', '.', ':', '+', '-'):
                if current:
                    words.append((current, 'normal'))
                    current = ""
                words.append((c, 'delimiter'))
            else:
                current += c
        char_idx += 1
        
    if current:
        words.append((current, 'string' if in_string else 'normal'))
        
    for val, kind in words:
        run = p.add_run()
        run.text = val
        if kind == 'comment':
            color = COLOR_TEXT_MUTED
            italic, bold = True, False
        elif kind == 'string':
            color = COLOR_ACCENT_GREEN
            italic, bold = False, False
        elif kind == 'normal' and val in keywords:
            color = COLOR_ACCENT_PINK
            italic, bold = False, True
        elif kind == 'normal' and val.isnumeric():
            color = COLOR_ACCENT_BLUE
            italic, bold = False, False
        elif kind == 'delimiter':
            color = COLOR_TEXT_SECONDARY
            italic, bold = False, False
        else:
            color = COLOR_TEXT_PRIMARY
            italic, bold = False, False
        format_run(run, font_name="Consolas", font_size=Pt(9.5), color=color, bold=bold, italic=italic)

def add_code_block(slide, left, top, width, height, code_str, is_bash=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(244, 244, 245)
    shape.line.color.rgb = COLOR_BORDER
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    lines = code_str.strip().split('\n')
    keywords = {'from', 'import', 'def', 'class', 'return', 'print', 'as'} if not is_bash else {'gcloud', 'vertex-ai', 'create', 'sync'}
    
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        
        if is_bash:
            if line.strip().startswith('#'):
                run = p.add_run()
                run.text = line
                format_run(run, font_name="Consolas", font_size=Pt(9.5), color=COLOR_TEXT_MUTED, italic=True)
            else:
                # Basic highlight for flags and command names
                words = line.split(' ')
                for w_idx, w in enumerate(words):
                    run = p.add_run()
                    run.text = w + (" " if w_idx < len(words)-1 else "")
                    if w.startswith('--'):
                        color = COLOR_ACCENT_PINK
                    elif w in keywords:
                        color = COLOR_ACCENT_BLUE
                    else:
                        color = COLOR_TEXT_PRIMARY
                    format_run(run, font_name="Consolas", font_size=Pt(9.5), color=color)
        else:
            parse_code_line(p, line, keywords)


# ==================== SLIDE 1: TITLE SLIDE ====================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

# Main Title & Subtitle in single textbox
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.2))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

p_cat = tf.paragraphs[0]
p_cat.text = "GOOGLE CLOUD MLOPS"
p_cat.space_after = Pt(10)
format_run(p_cat.runs[0], font_name="Segoe UI", font_size=Pt(13), color=COLOR_ACCENT_BLUE, bold=True)

p_title = tf.add_paragraph()
p_title.text = "Vertex AI Feature Store"
p_title.space_after = Pt(16)
format_run(p_title.runs[0], font_name="Segoe UI", font_size=Pt(48), color=COLOR_TEXT_PRIMARY, bold=True)

p_sub = tf.add_paragraph()
p_sub.text = "Unifying ML Feature Engineering, Low-Latency Serving, and Point-in-Time Correct Offline Datasets"
p_sub.space_after = Pt(36)
format_run(p_sub.runs[0], font_name="Segoe UI", font_size=Pt(16), color=COLOR_TEXT_SECONDARY)

# Presenter Info
pres_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(0.5))
tf_pres = pres_box.text_frame
tf_pres.word_wrap = True
tf_pres.margin_left = tf_pres.margin_right = tf_pres.margin_top = tf_pres.margin_bottom = Inches(0)
p_pres = tf_pres.paragraphs[0]
p_pres.text = "Presenter: Bethanasamy Rajamani  |  Audience: Professional ML Engineer Candidates"
format_run(p_pres.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_MUTED)


# ==================== SLIDE 2: THE CORE MLOPS CHALLENGES ====================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_slide_header(slide2, "The Problem", "Core Challenges in Production ML")

bullets2 = [
    ("Train-Serve Skew", "ML models fail when prediction features are calculated differently than training features."),
    ("Data Leakage", "Incorporating \"future data\" into training sets during joins, invalidating offline evaluation validation."),
    ("Redundant Engineering", "Separate data engineering teams re-calculating the exact same feature values, driving up storage/compute costs.")
]
add_bullets(slide2, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets2)
add_example_box(slide2, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "A credit card fraud detector checks transaction velocities. If training calculated velocities using daily averages, but online inference queries them using 5-minute windows, the mismatch (skew) causes false negatives.")

tf_card2 = add_card_box(slide2, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Traditional Databases Fall Short")
p_desc2 = tf_card2.add_paragraph()
p_desc2.text = "Relational databases (SQL) and analytical warehouses (BigQuery) cannot serve features with sub-15ms latency at scale. Conversely, transactional databases lack point-in-time snapshot mechanics, creating massive data leakage risks during dataset exports."
p_desc2.space_after = Pt(16)
format_run(p_desc2.runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_SECONDARY)

p_ex2 = tf_card2.add_paragraph()
p_ex2.text = "Example: Both fraud and billing teams writing two separate Spark jobs to calculate daily_user_active_hours leads to code redundancy and cost waste."
format_run(p_ex2.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_ACCENT_BLUE, bold=True)


# ==================== SLIDE 3: WHAT IS VERTEX AI FEATURE STORE? ====================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_slide_header(slide3, "Architectural Overview", "What is Vertex AI Feature Store?")

# Left intro text
intro_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.8), Inches(0.5))
p_intro = intro_box.text_frame.paragraphs[0]
p_intro.text = "A centralized, fully-managed Google Cloud registry to share, discover, and serve ML features."
format_run(p_intro.runs[0], font_name="Segoe UI", font_size=Pt(13), color=COLOR_TEXT_PRIMARY, bold=True)

bullets3 = [
    ("Single Source of Truth", "Unifies preprocessing pipelines across batch and real-time."),
    ("Serverless Scale", "No infrastructure management. Scaling Bigtable and BigQuery slots are handled entirely by Google Cloud."),
    ("Integrations", "Built-in lineage hooks into Vertex AI Model Registry, Pipelines, and Model Monitoring.")
]
add_bullets(slide3, Inches(0.8), Inches(2.1), Inches(6.8), Inches(2.9), bullets3)
add_example_box(slide3, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "A ride-sharing app uses online serving to retrieve driver ratings in 8ms during matching. For historical pricing analysis, it uses offline serving to export monthly coordinates.")

tf_card3 = add_card_box(slide3, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Operational Benefits")

# Benefit 1
p1 = tf_card3.add_paragraph()
p1.text = "⚡ Sub-15ms Online Serving Latency"
p1.space_before = Pt(4)
format_run(p1.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_ACCENT_BLUE, bold=True)
p1_desc = tf_card3.add_paragraph()
p1_desc.text = "Powered by Cloud Bigtable key-value storage, delivering high-throughput gRPC lookups for real-time model inference at scale."
p1_desc.space_after = Pt(10)
format_run(p1_desc.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_TEXT_SECONDARY)

# Benefit 2
p2 = tf_card3.add_paragraph()
p2.text = "🔄 Zero-Copy BigQuery Views Integration"
format_run(p2.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_ACCENT_GREEN, bold=True)
p2_desc = tf_card3.add_paragraph()
p2_desc.text = "Queries BigQuery tables or logical views directly without copying physical data, eliminating duplication storage costs and pipeline lags."
p2_desc.space_after = Pt(10)
format_run(p2_desc.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_TEXT_SECONDARY)

# Benefit 3
p3 = tf_card3.add_paragraph()
p3.text = "🔗 Fully Auditable Metadata Lineage"
format_run(p3.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_ACCENT_PINK, bold=True)
p3_desc = tf_card3.add_paragraph()
p3_desc.text = "Integrated with Vertex ML Metadata to track feature origin, transformations, registry audits, and downstream model consumers."
format_run(p3_desc.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 4: DATA MODELING HIERARCHY ====================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_slide_header(slide4, "Data Modeling", "Vertex AI Feature Store Hierarchy")

bullets4 = [
    ("Entity Type", "Represents a primary key tracking data records (e.g. user_id, product_id)."),
    ("Feature Group", "Points directly to a source table/view in BigQuery containing attributes."),
    ("Feature View", "A read-only representation of features. Synchronizes directly with BigQuery without copying raw data.")
]
add_bullets(slide4, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets4)
add_example_box(slide4, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "Entity Type = customer_id. Feature Group = customer_orders (linked to BigQuery). Feature View = checkout_view (exposing features like last_order_amount).")

tf_card4 = add_card_box(slide4, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Key Design Pattern")
p_desc4 = tf_card4.add_paragraph()
p_desc4.text = "You define Feature Groups to map BigQuery schemas. You register individual features within that group. Feature Views then expose specific subsets of these features to online serving APIs."
format_run(p_desc4.runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 5: BATCH & STREAMING INGESTIONS ====================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_slide_header(slide5, "Ingestion Pipelines", "Batch & Streaming Ingestions")

bullets5 = [
    ("Batch Ingestion (BigQuery)", "Scheduled sync jobs that pull pre-calculated features. Animate packets showing scheduled sync runs."),
    ("Streaming Ingestion (Pub/Sub)", "Ingests features in real-time. Events flow through Pub/Sub and Dataflow directly to the Feature Store.")
]
add_bullets(slide5, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets5)
add_example_box(slide5, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "Streaming Ingestion writes real-time customer clicks from Pub/Sub to the online store instantly. Batch ingestion synchronizes user credit ratings from BigQuery every midnight.")

# Draw Ingestion Diagram on the right
add_card_box(slide5, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Active Data Ingestion Pipeline")

# Nodes (shapes)
# 1. BigQuery (Batch)
bq_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.6), Inches(3.5), Inches(0.7))
bq_shape.fill.solid()
bq_shape.fill.fore_color.rgb = COLOR_CARD
bq_shape.line.color.rgb = COLOR_ACCENT_PINK
bq_shape.line.width = Pt(1.5)
tf_bq = bq_shape.text_frame
tf_bq.word_wrap = True
tf_bq.paragraphs[0].text = "BigQuery\n(Batch Source)"
tf_bq.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_bq.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_PRIMARY, bold=True)

# 2. Pub/Sub (Streaming)
ps_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.4), Inches(3.5), Inches(0.7))
ps_shape.fill.solid()
ps_shape.fill.fore_color.rgb = COLOR_CARD
ps_shape.line.color.rgb = COLOR_ACCENT_GREEN
ps_shape.line.width = Pt(1.5)
tf_ps = ps_shape.text_frame
tf_ps.word_wrap = True
tf_ps.paragraphs[0].text = "Pub/Sub\n(Streaming Source)"
ps_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_ps.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_PRIMARY, bold=True)

# 3. Feature Store (Center)
fs_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(4.0), Inches(3.5), Inches(0.7))
fs_shape.fill.solid()
fs_shape.fill.fore_color.rgb = COLOR_CARD
fs_shape.line.color.rgb = COLOR_ACCENT_BLUE
fs_shape.line.width = Pt(2)
tf_fs = fs_shape.text_frame
tf_fs.word_wrap = True
tf_fs.paragraphs[0].text = "Feature Store\n(Active Registry)"
fs_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_fs.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_PRIMARY, bold=True)

# Connector arrows
draw_arrow_line(slide5, Inches(10.25), Inches(3.3), Inches(10.25), Inches(4.0), COLOR_ACCENT_PINK)
draw_arrow_line(slide5, Inches(10.25), Inches(5.4), Inches(10.25), Inches(4.7), COLOR_ACCENT_GREEN)


# ==================== SLIDE 6: ONLINE SERVING ARCHITECTURE ====================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_slide_header(slide6, "Serving (Low Latency)", "Online Serving Mechanics")

bullets6 = [
    ("Cloud Bigtable", "The default serving engine. Scaled dynamically by Google Cloud to support massive, concurrent online requests."),
    ("Read Feature Values API", "Online endpoints request feature values by passing the Entity ID (e.g. user_id) and receiving vectors in milliseconds."),
    ("Zero-Latency Cache", "Feature Store buffers features to prevent server overloads during web traffic spikes.")
]
add_bullets(slide6, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets6)
add_example_box(slide6, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "The inference server queries user_1984 on port 8080 and fetches the feature vector [0.85 (ctr), 'Premium' (tier)] from Cloud Bigtable in 6.8ms.")

# Draw serving diagram
add_card_box(slide6, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Online serving gRPC query loop")

# Node: Model API (Circle)
api_shape = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(3.4), Inches(1.2), Inches(1.2))
api_shape.fill.solid()
api_shape.fill.fore_color.rgb = COLOR_CARD
api_shape.line.color.rgb = COLOR_TEXT_SECONDARY
api_shape.line.width = Pt(1.5)
tf_api = api_shape.text_frame
tf_api.paragraphs[0].text = "Model API"
api_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_api.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# Node: Bigtable (Circle)
bt_shape = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(3.4), Inches(1.2), Inches(1.2))
bt_shape.fill.solid()
bt_shape.fill.fore_color.rgb = COLOR_CARD
bt_shape.line.color.rgb = COLOR_ACCENT_GREEN
bt_shape.line.width = Pt(2)
tf_bt = bt_shape.text_frame
tf_bt.paragraphs[0].text = "Bigtable"
bt_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_bt.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# Connecting Arrows (Request & Response)
draw_arrow_line(slide6, Inches(9.7), Inches(3.7), Inches(10.8), Inches(3.7), COLOR_ACCENT_GREEN)
draw_arrow_line(slide6, Inches(10.8), Inches(4.3), Inches(9.7), Inches(4.3), COLOR_ACCENT_BLUE)

# Latency badge
badge = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.65), Inches(2.6), Inches(1.2), Inches(0.5))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(244, 244, 245)
badge.line.color.rgb = COLOR_BORDER
badge.line.width = Pt(1)
tf_badge = badge.text_frame
tf_badge.paragraphs[0].text = "6.8ms"
badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_badge.paragraphs[0].runs[0], font_name="Consolas", font_size=Pt(11.5), color=COLOR_ACCENT_GREEN, bold=True)


# ==================== SLIDE 7: OFFLINE SERVING & HISTORICAL JOINS ====================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_slide_header(slide7, "Serving (Analytical)", "Offline Serving & point-in-Time Joins")

bullets7 = [
    ("Point-in-Time joins", "Prevents data leakage by retrieving feature values exactly as they existed at the timestamp of a historical training event."),
    ("Batch Exports", "Exports features to BigQuery or Cloud Storage folders to build massive training sets."),
    ("Scale", "Runs analytical joins across millions of training records, leveraging BigQuery slot allocations.")
]
add_bullets(slide7, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets7)
add_example_box(slide7, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "To train a churn model on a checkout event from Oct 5, 2:00 PM, the Feature Store runs a point-in-time join that retrieves features prior to 2:00 PM, ignoring any subsequent activity.")

# Draw timeline travel diagram
add_card_box(slide7, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Point-in-Time \"Time Travel\" Join")

# Timeline horizontal line
timeline = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(3.8), Inches(3.7), Inches(0.04))
timeline.fill.solid()
timeline.fill.fore_color.rgb = COLOR_BORDER
timeline.line.fill.background()

# Dots represent updates
dot1 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(3.7), Inches(0.25), Inches(0.25))
dot1.fill.solid()
dot1.fill.fore_color.rgb = COLOR_ACCENT_GREEN
dot1.line.fill.background()

dot2 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(3.7), Inches(0.25), Inches(0.25))
dot2.fill.solid()
dot2.fill.fore_color.rgb = COLOR_ACCENT_GREEN
dot2.line.fill.background()

dot3 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(3.7), Inches(0.25), Inches(0.25))
dot3.fill.solid()
dot3.fill.fore_color.rgb = COLOR_TEXT_MUTED
dot3.line.fill.background()

# Slider vertical line (T_event)
slider = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.7), Inches(2.6), Inches(0.04), Inches(2.0))
slider.fill.solid()
slider.fill.fore_color.rgb = COLOR_ACCENT_BLUE
slider.line.fill.background()

# Text Labels
lbl_box = slide7.shapes.add_textbox(Inches(8.4), Inches(2.1), Inches(3.7), Inches(0.4))
tf_lbl = lbl_box.text_frame
tf_lbl.word_wrap = True
tf_lbl.paragraphs[0].text = "Dynamic Join Timestamp (T_event)"
lbl_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_lbl.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_ACCENT_BLUE, bold=True)

lbl_warn = slide7.shapes.add_textbox(Inches(8.4), Inches(4.7), Inches(3.7), Inches(0.4))
tf_warn = lbl_warn.text_frame
tf_warn.word_wrap = True
tf_warn.paragraphs[0].text = "[Future Data Blocked]"
lbl_warn.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_warn.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_ACCENT_PINK, bold=True)


# ==================== SLIDE 8: VERTEX AI SDK (PYTHON) ====================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_slide_header(slide8, "Developer Guide", "Vertex AI SDK (Python)")

# Intro desc
tb_desc8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4))
p_desc8 = tb_desc8.text_frame.paragraphs[0]
p_desc8.text = "Interacting with Vertex AI Feature Store programmatically:"
format_run(p_desc8.runs[0], font_name="Segoe UI", font_size=Pt(13), color=COLOR_TEXT_PRIMARY)

code8 = """
from google.cloud import aiplatform

aiplatform.init(project="my-gcp-project", location="us-central1")

# Define reference to registered Feature View
feature_view = aiplatform.FeatureView(
    feature_view_name="user_features",
    feature_online_store_id="user_online_store"
)

# Fetch real-time online feature values
response = feature_view.read_feature_values(
    entity_ids=["user_10294"]
)
print(response.values)
"""
add_code_block(slide8, Inches(0.8), Inches(1.9), Inches(11.73), Inches(3.4), code8, is_bash=False)
add_example_box(slide8, Inches(0.8), Inches(5.5), Inches(11.73), Inches(1.2), 
    "This Python query is integrated inside a web-serving Flask/FastAPI app to dynamically fetch driver ratings and enrich matching payloads before passing variables to model predictors.")


# ==================== SLIDE 9: GOOGLE CLOUD CLI REFERENCE ====================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_slide_header(slide9, "Operations Guide", "Google Cloud CLI (gcloud) Reference")

# Intro desc
tb_desc9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4))
p_desc9 = tb_desc9.text_frame.paragraphs[0]
p_desc9.text = "Managing Feature Store infrastructure using terminal commands:"
format_run(p_desc9.runs[0], font_name="Segoe UI", font_size=Pt(13), color=COLOR_TEXT_PRIMARY)

code9 = """
# Create online store instance backed by Bigtable
gcloud vertex-ai feature-online-stores create user-online-store \
    --region=us-central1 \
    --bigtable-min-node-count=1 \
    --bigtable-max-node-count=5

# Synchronize feature view with BigQuery source data
gcloud vertex-ai feature-views sync user-features \
    --feature-online-store=user-online-store \
    --region=us-central1
"""
add_code_block(slide9, Inches(0.8), Inches(1.9), Inches(11.73), Inches(3.4), code9, is_bash=True)
add_example_box(slide9, Inches(0.8), Inches(5.5), Inches(11.73), Inches(1.2), 
    "MLOps engineers run these gcloud scripts inside Cloud Build templates to automatically create and update feature registries when new schemas are merged into main branches.")


# ==================== SLIDE 10: INTEGRATION WITH MODEL MONITORING ====================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_slide_header(slide10, "MLOps Monitoring", "Model Monitoring Integration")

bullets10 = [
    ("Continuous Skew Checks", "Automated alerts if serving feature distributions drift from baseline training datasets (TFRecord profiles)."),
    ("PSI Metric", "Model Monitoring calculates Population Stability Index. If PSI exceeds 0.20, alerts notify engineering teams."),
    ("Automated Retraining", "Integration with Pub/Sub allows alerting systems to launch retraining DAGs automatically.")
]
add_bullets(slide10, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets10)
add_example_box(slide10, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "Monitoring user age distributions. Baseline training mean was 24. Shifted average in serving hits 42 (PSI = 0.28). System alerts trigger an automated retraining pipeline.")

tf_card10 = add_card_box(slide10, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "MLOps Operations Loop")
p_desc10 = tf_card10.add_paragraph()
p_desc10.text = "Unifying features within the Feature Store ensures the baseline reference schemas remain identical, eliminating manual calculations during drift checks."
format_run(p_desc10.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 11: IAM & SECURITY GOVERNANCE ====================
slide11 = prs.slides.add_slide(blank_layout)
set_slide_background(slide11)
add_slide_header(slide11, "Security & Governance", "IAM Permissions & Access Controls")

bullets11 = [
    ("roles/aiplatform.admin", "Grants full control over Feature Store creation, schema modification, and resource deletions."),
    ("roles/aiplatform.user", "Allows listing, searching, reading, and exporting feature values. Default role for ML engineers and clients."),
    ("VPC Service Controls", "Lock Feature Store endpoints inside a private security boundary, protecting datasets from exfiltration.")
]
add_bullets(slide11, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets11)
add_example_box(slide11, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "The online prediction container's service account is granted the aiplatform.user role, enabling it to query features on port 8080 while blocking write configurations.")

tf_card11 = add_card_box(slide11, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Security Rule")
p_desc11 = tf_card11.add_paragraph()
p_desc11.text = "Service accounts attached to prediction serving containers must possess the roles/aiplatform.user role to successfully pull online features on port 8080."
format_run(p_desc11.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 12: ARCHITECT'S SUMMARY ====================
slide12 = prs.slides.add_slide(blank_layout)
set_slide_background(slide12)
add_slide_header(slide12, "Takeaways", "Architect's Checklist")

bullets12 = [
    ("Train-Serve Skew", "Eliminate skew by serving and training from the same Feature Store views."),
    ("Data Leakage", "Eliminate leakage by deploying point-in-time snapshot joins."),
    ("Online Store", "Backed by Bigtable for sub-15ms serving latencies."),
    ("Offline Store", "Integrates with BigQuery for large analytical dataset creation.")
]
add_bullets(slide12, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets12)
add_example_box(slide12, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "On the exam, if they request sub-10ms serving, historical data timelines, or data quality assurances for tabular features, Vertex AI Feature Store is the correct answer.")

tf_card12 = add_card_box(slide12, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Exam Tip", border_color=COLOR_ACCENT_GREEN)
p_desc12 = tf_card12.add_paragraph()
p_desc12.text = "If a scenario requires real-time serving latency, user sharing, or data leakage protection, **Vertex AI Feature Store** is always the correct GCP ML solution architecture."
format_run(p_desc12.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_SECONDARY)


# Save PowerPoint deck
output_path = "Vertex_AI_Feature_Store.pptx"
prs.save(output_path)
print(f"PowerPoint slide deck created successfully: {output_path}")
