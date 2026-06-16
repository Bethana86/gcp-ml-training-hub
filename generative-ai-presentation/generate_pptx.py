import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation and Widescreen Mode (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Constants
COLOR_BG = RGBColor(228, 228, 230)
COLOR_CARD = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(212, 212, 216)
COLOR_TEXT_PRIMARY = RGBColor(24, 24, 27)
COLOR_TEXT_SECONDARY = RGBColor(63, 63, 70)
COLOR_TEXT_MUTED = RGBColor(113, 113, 122)
COLOR_ACCENT_BLUE = RGBColor(29, 78, 216)
COLOR_ACCENT_GREEN = RGBColor(4, 120, 87)
COLOR_ACCENT_PINK = RGBColor(190, 24, 93)

# Blank layout helper
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def format_run(run, font_name="Segoe UI", font_size=Pt(14), color=COLOR_TEXT_SECONDARY, bold=False, italic=False):
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic

def add_slide_header(slide, category, title):
    # Category (small uppercase text)
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = Inches(0)
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    format_run(p_cat.runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_ACCENT_BLUE, bold=True)
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = Inches(0)
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    format_run(p_title.runs[0], font_name="Segoe UI", font_size=Pt(28), color=COLOR_TEXT_PRIMARY, bold=True)

def add_example_box(slide, left, top, width, height, example_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(244, 244, 245)
    shape.line.color.rgb = COLOR_ACCENT_BLUE
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    
    p0 = tf.paragraphs[0]
    p0.text = "REAL-WORLD EXAMPLE"
    p0.space_after = Pt(3)
    format_run(p0.runs[0], font_name="Segoe UI", font_size=Pt(9.5), color=COLOR_ACCENT_BLUE, bold=True)
    
    p1 = tf.add_paragraph()
    p1.text = example_text
    format_run(p1.runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_SECONDARY, italic=True)

def add_card_box(slide, left, top, width, height, card_title, border_color=COLOR_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    p0 = tf.paragraphs[0]
    p0.text = card_title
    p0.space_after = Pt(12)
    format_run(p0.runs[0], font_name="Segoe UI", font_size=Pt(16), color=COLOR_TEXT_PRIMARY, bold=True)
    
    return tf

def add_bullets(slide, left, top, width, height, bullets):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    
    for idx, (title, desc) in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        
        run_arrow = p.add_run()
        run_arrow.text = "→  "
        format_run(run_arrow, font_name="Segoe UI", font_size=Pt(13), color=COLOR_ACCENT_BLUE, bold=True)
        
        run_title = p.add_run()
        run_title.text = title + ": "
        format_run(run_title, font_name="Segoe UI", font_size=Pt(12.5), color=COLOR_TEXT_PRIMARY, bold=True)
        
        run_desc = p.add_run()
        run_desc.text = desc
        format_run(run_desc, font_name="Segoe UI", font_size=Pt(12.5), color=COLOR_TEXT_SECONDARY)

def draw_arrow_line(slide, x1, y1, x2, y2, color):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml import parse_xml
    
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    
    # Add arrowhead to tailEnd (end of path)
    ln = connector.line._get_or_add_ln()
    tailEnd = parse_xml(
        '<a:tailEnd type="arrow" w="med" len="med" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    )
    ln.append(tailEnd)
    return connector

def parse_code_line(p, line, keywords):
    # Simple inline parser for python formatting
    char_idx = 0
    words = []
    current = ""
    in_string = False
    string_char = None
    
    while char_idx < len(line):
        c = line[char_idx]
        
        if c == '#' and not in_string:
            if current:
                words.append((current, 'normal'))
                current = ""
            words.append((line[char_idx:], 'comment'))
            break
            
        if c in ('"', "'") and (char_idx == 0 or line[char_idx-1] != '\\'):
            if in_string:
                if c == string_char:
                    current += c
                    words.append((current, 'string'))
                    current = ""
                    in_string = False
                    string_char = None
                else:
                    current += c
            else:
                if current:
                    words.append((current, 'normal'))
                    current = ""
                in_string = True
                string_char = c
                current += c
            char_idx += 1
            continue
            
        if in_string:
            current += c
        else:
            if c in (' ', '(', ')', '[', ']', '{', '}', '=', ',', '.', ':', '+', '-'):
                if current:
                    words.append((current, 'normal'))
                    current = ""
                words.append((c, 'delimiter'))
            else:
                current += c
        char_idx += 1
        
    if current:
        words.append((current, 'string' if in_string else 'normal'))
        
    for val, kind in words:
        run = p.add_run()
        run.text = val
        if kind == 'comment':
            color = COLOR_TEXT_MUTED
            italic, bold = True, False
        elif kind == 'string':
            color = COLOR_ACCENT_GREEN
            italic, bold = False, False
        elif kind == 'normal' and val in keywords:
            color = COLOR_ACCENT_PINK
            italic, bold = False, True
        elif kind == 'normal' and val.isnumeric():
            color = COLOR_ACCENT_BLUE
            italic, bold = False, False
        elif kind == 'delimiter':
            color = COLOR_TEXT_SECONDARY
            italic, bold = False, False
        else:
            color = COLOR_TEXT_PRIMARY
            italic, bold = False, False
        format_run(run, font_name="Consolas", font_size=Pt(9.5), color=color, bold=bold, italic=italic)

def add_code_block(slide, left, top, width, height, code_str, is_bash=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(244, 244, 245)
    shape.line.color.rgb = COLOR_BORDER
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    lines = code_str.strip().split('\n')
    keywords = {'from', 'import', 'def', 'class', 'return', 'print', 'as'} if not is_bash else {'curl', 'gcloud', 'auth', 'print-access-token'}
    
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        
        if is_bash:
            if line.strip().startswith('#'):
                run = p.add_run()
                run.text = line
                format_run(run, font_name="Consolas", font_size=Pt(9.5), color=COLOR_TEXT_MUTED, italic=True)
            else:
                # Basic highlight for flags and command names
                words = line.split(' ')
                for w_idx, w in enumerate(words):
                    run = p.add_run()
                    run.text = w + (" " if w_idx < len(words)-1 else "")
                    if w.startswith('--') or w.startswith('-H') or w.startswith('-X') or w.startswith('-d'):
                        color = COLOR_ACCENT_PINK
                    elif w in keywords:
                        color = COLOR_ACCENT_BLUE
                    else:
                        color = COLOR_TEXT_PRIMARY
                    format_run(run, font_name="Consolas", font_size=Pt(9.5), color=color)
        else:
            parse_code_line(p, line, keywords)

def add_qa_slide(category, title, question_text, options, correct_letter, explanation, ref_link):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_header(slide, category, title)
    
    # Add Question Text and Options on the left
    tb_q = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.8), Inches(5.2))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = Inches(0)
    
    p_q = tf_q.paragraphs[0]
    p_q.text = question_text
    p_q.space_after = Pt(14)
    format_run(p_q.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_TEXT_PRIMARY, bold=True)
    
    for opt_letter, opt_text in options:
        p_opt = tf_q.add_paragraph()
        p_opt.space_after = Pt(8)
        
        run_let = p_opt.add_run()
        run_let.text = f"{opt_letter}.  "
        format_run(run_let, font_name="Segoe UI", font_size=Pt(11), color=COLOR_ACCENT_BLUE, bold=True)
        
        run_txt = p_opt.add_run()
        run_txt.text = opt_text
        format_run(run_txt, font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_SECONDARY)
        
    # Add Answer Card on the right
    tf_card = add_card_box(slide, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Validation Panel")
    
    p_ans = tf_card.add_paragraph()
    p_ans.text = f"Correct Answer: {correct_letter}"
    p_ans.space_after = Pt(10)
    format_run(p_ans.runs[0], font_name="Segoe UI", font_size=Pt(14), color=COLOR_ACCENT_GREEN, bold=True)
    
    p_exp = tf_card.add_paragraph()
    p_exp.text = f"Explanation: {explanation}"
    p_exp.space_after = Pt(14)
    format_run(p_exp.runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_SECONDARY)
    
    p_ref = tf_card.add_paragraph()
    p_ref.text = f"GCP Reference: {ref_link}"
    format_run(p_ref.runs[0], font_name="Segoe UI", font_size=Pt(9.5), color=COLOR_TEXT_MUTED, italic=True)


# ==================== SLIDE 1: TITLE SLIDE ====================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.2))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

p_cat = tf.paragraphs[0]
p_cat.text = "GOOGLE CLOUD GENAI"
p_cat.space_after = Pt(10)
format_run(p_cat.runs[0], font_name="Segoe UI", font_size=Pt(13), color=COLOR_ACCENT_BLUE, bold=True)

p_title = tf.add_paragraph()
p_title.text = "Introduction to Generative AI"
p_title.space_after = Pt(16)
format_run(p_title.runs[0], font_name="Segoe UI", font_size=Pt(48), color=COLOR_TEXT_PRIMARY, bold=True)

p_sub = tf.add_paragraph()
p_sub.text = "Core Concepts, Transformer Architectures, LLM Pre-training, and Enterprise Prompt Engineering"
p_sub.space_after = Pt(36)
format_run(p_sub.runs[0], font_name="Segoe UI", font_size=Pt(16), color=COLOR_TEXT_SECONDARY)

# Presenter Info
pres_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(0.5))
tf_pres = pres_box.text_frame
tf_pres.word_wrap = True
tf_pres.margin_left = tf_pres.margin_right = tf_pres.margin_top = tf_pres.margin_bottom = Inches(0)
p_pres = tf_pres.paragraphs[0]
p_pres.text = "Presenter: Bethanasamy Rajamani  |  Audience: Professional ML Engineer Candidates"
format_run(p_pres.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_MUTED)


# ==================== SLIDE 2: CORE CHALLENGES IN GENAI APPLICATIONS ====================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_slide_header(slide2, "The Problem", "Core Challenges in GenAI Applications")

bullets2 = [
    ("Hallucination Risks", "Models generate fluent but factually incorrect statements due to statistical next-token optimization."),
    ("Context Window Limits", "Transformers can process only a limited number of tokens in memory, leading to retrieval failures on massive corpora."),
    ("High Inference Costs", "Auto-regressive decoding is memory-bound, demanding persistent high-tier GPU/TPU resources.")
]
add_bullets(slide2, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets2)
add_example_box(slide2, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "A customer support chatbot attempts to quote insurance terms. Without grounding, it fabricates a refund policy (hallucination). At 10,000 queries per minute, standard API serving costs scale exponentially (high inference costs).")

tf_card2 = add_card_box(slide2, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Why Standard NLP Falls Short")
p_desc2 = tf_card2.add_paragraph()
p_desc2.text = "Traditional rule-based and classification NLP models (such as BERT or simple LSTM classifiers) are excellent for sentiment scoring or named entity extraction. However, they lack the semantic reasoning, zero-shot generalization, and synthesis capabilities required to build interactive reasoning agents."
p_desc2.space_after = Pt(16)
format_run(p_desc2.runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_SECONDARY)

p_ex2 = tf_card2.add_paragraph()
p_ex2.text = "Example: Writing custom classification rules for 50 different document layouts is inefficient compared to prompting a foundation model to extract structured JSON."
format_run(p_ex2.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_ACCENT_BLUE, bold=True)


# ==================== SLIDE 3: WHAT IS GENERATIVE AI? ====================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_slide_header(slide3, "Architectural Overview", "What is Generative AI?")

intro_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.8), Inches(0.5))
p_intro = intro_box.text_frame.paragraphs[0]
p_intro.text = "A category of artificial intelligence that creates new content, utilizing massive pre-trained transformer weights."
format_run(p_intro.runs[0], font_name="Segoe UI", font_size=Pt(13), color=COLOR_TEXT_PRIMARY, bold=True)

bullets3 = [
    ("Centralized Foundation Models", "Trained on internet-scale datasets, forming baseline reasoning capabilities."),
    ("Multimodality", "Natively processes and connects diverse modalities including text, images, video, and audio code."),
    ("Transfer Learning", "Repurposes general knowledge to specific downstream tasks via fine-tuning or in-context prompting.")
]
add_bullets(slide3, Inches(0.8), Inches(2.1), Inches(6.8), Inches(2.9), bullets3)
add_example_box(slide3, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "A business analyst inputs natural language questions. The model generates both the SQL query and an explanatory summary, translating unstructured language into structured system database queries.")

tf_card3 = add_card_box(slide3, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Operational Benefits")

# Benefit 1
p1 = tf_card3.add_paragraph()
p1.text = "⚡ Zero-Shot Reasoning Capabilities"
p1.space_before = Pt(4)
format_run(p1.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_ACCENT_BLUE, bold=True)
p1_desc = tf_card3.add_paragraph()
p1_desc.text = "Solves new, unseen tasks simply by receiving descriptive prompts without requiring model gradient updates or parameter fine-tuning."
p1_desc.space_after = Pt(10)
format_run(p1_desc.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_TEXT_SECONDARY)

# Benefit 2
p2 = tf_card3.add_paragraph()
p2.text = "🔄 API Scalability & Managed Access"
format_run(p2.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_ACCENT_GREEN, bold=True)
p2_desc = tf_card3.add_paragraph()
p2_desc.text = "Deploy models serverless via the Gemini API, eliminating custom hardware hosting, scaling, and load-balancer operations."
p2_desc.space_after = Pt(10)
format_run(p2_desc.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_TEXT_SECONDARY)

# Benefit 3
p3 = tf_card3.add_paragraph()
p3.text = "🔗 Rapid Application Deployment Loops"
format_run(p3.runs[0], font_name="Segoe UI", font_size=Pt(12), color=COLOR_ACCENT_PINK, bold=True)
p3_desc = tf_card3.add_paragraph()
p3_desc.text = "Decouples model engineering from software features. Application logic adjusts by updating system prompts rather than retraining pipelines."
format_run(p3_desc.runs[0], font_name="Segoe UI", font_size=Pt(10.5), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 4: TRANSFORMER ARCHITECTURE HIERARCHY ====================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_slide_header(slide4, "Model Architecture", "Transformer Architecture Hierarchy")

bullets4 = [
    ("Self-Attention Layer", "Calculates dynamic mathematical weights representing how tokens relate within a context window."),
    ("Encoder-Decoder Blocks", "Encodes input text into vector spaces, and decodes vectors to target sequences (ideal for translation)."),
    ("Decoder-Only Autoregressive", "Generates next tokens sequentially based on past contexts (backbone of modern Gemini/LLMs).")
]
add_bullets(slide4, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets4)
add_example_box(slide4, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "Input Tokens = [\"Generative\", \"AI\", \"is\"]. Self-attention maps that \"AI\" strongly relates to \"Generative\", predicting the next probable token \"transformational\".")

tf_card4 = add_card_box(slide4, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Key Design Pattern")
p_desc4 = tf_card4.add_paragraph()
p_desc4.text = "Raw inputs are tokenized and mapped to continuous vector Embeddings. Positional Encodings are added to maintain sequence order. Multi-head attention blocks then capture relationships in parallel, outputting normalized token probabilities."
format_run(p_desc4.runs[0], font_name="Segoe UI", font_size=Pt(11), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 5: PRE-TRAINING & FINE-TUNING PIPELINES ====================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_slide_header(slide5, "Ingestion Pipelines", "Pre-training & Fine-Tuning Pipelines")

bullets5 = [
    ("Unsupervised Pre-training", "Large-scale model training on raw text corpus to predict next words, building baseline semantics."),
    ("Supervised Fine-Tuning (SFT)", "Trains the base model on high-quality instruction-response pairs to learn task compliance."),
    ("RLHF Alignment", "Integrates Reinforcement Learning from Human Feedback to align outputs with safety and preferences.")
]
add_bullets(slide5, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets5)
add_example_box(slide5, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "Raw internet scrap builds a Base Model. Supervised training on 50,000 chat logs creates an SFT Model. Reward optimization aligned with human reviewers produces a final chat-ready LLM.")

# Draw Training Flow Diagram on the right
add_card_box(slide5, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Active Model Training Flow")

# Nodes (shapes)
# 1. Raw Text
t1_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(2.5), Inches(1.7), Inches(0.6))
t1_shape.fill.solid()
t1_shape.fill.fore_color.rgb = COLOR_CARD
t1_shape.line.color.rgb = COLOR_ACCENT_PINK
t1_shape.line.width = Pt(1.5)
tf_t1 = t1_shape.text_frame
tf_t1.word_wrap = True
tf_t1.paragraphs[0].text = "Raw Text\n(Unlabeled)"
tf_t1.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_t1.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# 2. Base Model
t2_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(2.5), Inches(1.7), Inches(0.6))
t2_shape.fill.solid()
t2_shape.fill.fore_color.rgb = COLOR_CARD
t2_shape.line.color.rgb = COLOR_ACCENT_PINK
t2_shape.line.width = Pt(1.5)
tf_t2 = t2_shape.text_frame
tf_t2.word_wrap = True
tf_t2.paragraphs[0].text = "Base Model\n(Weights)"
t2_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_t2.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# 3. SFT Model
t3_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(4.5), Inches(1.7), Inches(0.6))
t3_shape.fill.solid()
t3_shape.fill.fore_color.rgb = COLOR_CARD
t3_shape.line.color.rgb = COLOR_ACCENT_GREEN
t3_shape.line.width = Pt(1.5)
tf_t3 = t3_shape.text_frame
tf_t3.word_wrap = True
tf_t3.paragraphs[0].text = "SFT Model\n(Instruction Tuned)"
t3_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_t3.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# 4. RLHF Model
t4_shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(4.5), Inches(1.7), Inches(0.6))
t4_shape.fill.solid()
t4_shape.fill.fore_color.rgb = COLOR_CARD
t4_shape.line.color.rgb = COLOR_ACCENT_BLUE
t4_shape.line.width = Pt(2)
tf_t4 = t4_shape.text_frame
tf_t4.word_wrap = True
tf_t4.paragraphs[0].text = "RLHF Model\n(Aligned Preferences)"
t4_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_t4.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# Connector arrows
draw_arrow_line(slide5, Inches(10.0), Inches(2.8), Inches(10.5), Inches(2.8), COLOR_ACCENT_PINK)
draw_arrow_line(slide5, Inches(11.35), Inches(3.1), Inches(9.15), Inches(4.5), COLOR_ACCENT_GREEN)
draw_arrow_line(slide5, Inches(10.0), Inches(4.8), Inches(10.5), Inches(4.8), COLOR_ACCENT_BLUE)


# ==================== SLIDE 6: LLM INFERENCE MECHANICS ====================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_slide_header(slide6, "Serving (Low Latency)", "LLM Inference Mechanics")

bullets6 = [
    ("Tokenization & Embeddings", "Input text is broken down into unique numerical token IDs, mapped to semantic vector spaces."),
    ("KV Caching", "Stores key-value attention representations of past tokens, skipping redundant matrix calculations during generation loops."),
    ("Decoding Parameters", "Parameters like Temperature and Top-K/Top-P shape output selection, controlling creativity vs consistency.")
]
add_bullets(slide6, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets6)
add_example_box(slide6, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "An API endpoint receives a prompt. It tokenizes inputs, processes them through transformer layers in 24ms (Time-To-First-Token), and streams generated tokens sequentially over HTTP/2.")

# Draw serving diagram
add_card_box(slide6, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Online Token generation & KV Loop")

# Node: Model API (Circle)
api_shape = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(3.4), Inches(1.2), Inches(1.2))
api_shape.fill.solid()
api_shape.fill.fore_color.rgb = COLOR_CARD
api_shape.line.color.rgb = COLOR_TEXT_SECONDARY
api_shape.line.width = Pt(1.5)
tf_api = api_shape.text_frame
tf_api.paragraphs[0].text = "Model API"
api_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_api.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# Node: LLM Core (Circle)
bt_shape = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(3.4), Inches(1.2), Inches(1.2))
bt_shape.fill.solid()
bt_shape.fill.fore_color.rgb = COLOR_CARD
bt_shape.line.color.rgb = COLOR_ACCENT_GREEN
bt_shape.line.width = Pt(2)
tf_bt = bt_shape.text_frame
tf_bt.paragraphs[0].text = "LLM Core"
bt_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_bt.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_TEXT_PRIMARY, bold=True)

# Connecting Arrows (Request & Response)
draw_arrow_line(slide6, Inches(9.7), Inches(3.7), Inches(10.8), Inches(3.7), COLOR_ACCENT_GREEN)
draw_arrow_line(slide6, Inches(10.8), Inches(4.3), Inches(9.7), Inches(4.3), COLOR_ACCENT_BLUE)

# Latency badge
badge = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.65), Inches(2.6), Inches(1.2), Inches(0.5))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(244, 244, 245)
badge.line.color.rgb = COLOR_BORDER
badge.line.width = Pt(1)
tf_badge = badge.text_frame
tf_badge.paragraphs[0].text = "24ms"
badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_badge.paragraphs[0].runs[0], font_name="Consolas", font_size=Pt(11.5), color=COLOR_ACCENT_GREEN, bold=True)


# ==================== SLIDE 7: RETRIEVAL-AUGMENTED GENERATION (RAG) ====================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_slide_header(slide7, "Serving (Analytical)", "Retrieval-Augmented Generation (RAG)")

bullets7 = [
    ("Document Chunking", "Divides large enterprise manuals into overlapping semantic text snippets for indexing."),
    ("Vector Embeddings", "Encodes document chunks into dense coordinate vectors using text embedding models."),
    ("Semantic Join", "Searches vector spaces using cosine similarity to find document context segments matching user queries.")
]
add_bullets(slide7, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets7)
add_example_box(slide7, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "A user asks: \"What is our Q3 cloud budget?\" The system searches Vector Search, retrieves a table chunk from q3_financials.pdf, and injects it to ground the LLM's response.")

# Draw RAG search sweep diagram
add_card_box(slide7, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Query Vector Sweep & Semantic Match")

# Timeline horizontal line
timeline = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(3.8), Inches(3.7), Inches(0.04))
timeline.fill.solid()
timeline.fill.fore_color.rgb = COLOR_BORDER
timeline.line.fill.background()

# Dots represent document chunks in vector space
dot1 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(3.7), Inches(0.25), Inches(0.25))
dot1.fill.solid()
dot1.fill.fore_color.rgb = COLOR_ACCENT_GREEN
dot1.line.fill.background()

dot2 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(3.7), Inches(0.25), Inches(0.25))
dot2.fill.solid()
dot2.fill.fore_color.rgb = COLOR_ACCENT_GREEN
dot2.line.fill.background()

dot3 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(3.7), Inches(0.25), Inches(0.25))
dot3.fill.solid()
dot3.fill.fore_color.rgb = COLOR_TEXT_MUTED
dot3.line.fill.background()

# Slider vertical line (Query vector sweep)
slider = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.7), Inches(2.6), Inches(0.04), Inches(2.0))
slider.fill.solid()
slider.fill.fore_color.rgb = COLOR_ACCENT_BLUE
slider.line.fill.background()

# Text Labels
lbl_box = slide7.shapes.add_textbox(Inches(8.4), Inches(2.1), Inches(3.7), Inches(0.4))
tf_lbl = lbl_box.text_frame
tf_lbl.word_wrap = True
tf_lbl.paragraphs[0].text = "Query Vector Sweep (Cosine Lookup)"
lbl_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_lbl.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_ACCENT_BLUE, bold=True)

lbl_warn = slide7.shapes.add_textbox(Inches(8.4), Inches(4.7), Inches(3.7), Inches(0.4))
tf_warn = lbl_warn.text_frame
tf_warn.word_wrap = True
tf_warn.paragraphs[0].text = "[Semantic Context Injected]"
lbl_warn.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
format_run(tf_warn.paragraphs[0].runs[0], font_name="Segoe UI", font_size=Pt(10), color=COLOR_ACCENT_PINK, bold=True)


# ==================== SLIDE 8: LLM SAFETY & MONITORING INTEGRATION ====================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_slide_header(slide8, "MLOps Monitoring", "LLM Safety & Monitoring Integration")

bullets8 = [
    ("Content Moderation Filters", "Built-in safety hooks that evaluate model queries and responses for harassment, hate speech, and explicit blocks."),
    ("Semantic Drift Checks", "Monitoring client prompt query patterns over time to capture task drift or adversarial injection patterns."),
    ("Guardrails Integration", "Wrapping model routes with external middleware to filter toxic language prior to client presentation.")
]
add_bullets(slide8, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets8)
add_example_box(slide8, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "An input prompt contains toxic patterns. The Vertex AI safety layer intercepts the request, blocks execution, and returns a structured response indicating a blocked classification.")

tf_card8 = add_card_box(slide8, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "MLOps Operations Loop")
p_desc8 = tf_card8.add_paragraph()
p_desc8.text = "Deploying models in production requires active validation loops. Integrating safety checks in the prediction pipeline prevents reputation risks, toxic outputs, and jailbreak exploits."
format_run(p_desc8.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_SECONDARY)


# ==================== SLIDE 9: ARCHITECT'S SUMMARY ====================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_slide_header(slide9, "Takeaways", "Architect's Checklist")

bullets9 = [
    ("Zero-Shot/Few-Shot", "First attempt prompting configurations before investing in custom parameter tuning."),
    ("Retrieval-Augmented Gen", "Ground models using Vector Search to mitigate hallucinations and present fresh datasets."),
    ("Cost/Latency", "Select Gemini Flash for quick, budget serving. Reserve Pro for complex reasoning."),
    ("Governance", "Enforce IAM role boundaries and safety filter limits at the endpoint level.")
]
add_bullets(slide9, Inches(0.8), Inches(1.8), Inches(6.8), Inches(3.2), bullets9)
add_example_box(slide9, Inches(0.8), Inches(5.2), Inches(6.8), Inches(1.5), 
    "On the exam, if a scenario asks to connect private enterprise PDF documents with LLM outputs without model retraining, the answer is always RAG with Vector Search.")

tf_card9 = add_card_box(slide9, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9), "Exam Tip", border_color=COLOR_ACCENT_GREEN)
p_desc9 = tf_card9.add_paragraph()
p_desc9.text = "If the exam asks to adapt a model to proprietary, frequently changing documentation under tight budgets, **Retrieval-Augmented Generation (RAG)** is the correct choice over custom fine-tuning."
format_run(p_desc9.runs[0], font_name="Segoe UI", font_size=Pt(11.5), color=COLOR_TEXT_SECONDARY)


# ==================== PRACTICE SCENARIOS (Q&A SLIDES 10-19) ====================
scenarios = [
    {
        "category": "Practice Scenario 1 of 10",
        "title": "Choosing Fine-Tuning vs RAG",
        "question": "A financial service company wants to build an internal virtual assistant to answer questions about dynamic company policies and daily updated mortgage rates. The documents contain sensitive proprietary information. What is the most cost-effective and architecturally sound pattern on GCP to implement this without exposing the base model to stale information?",
        "options": [
            ("A", "Run a Vertex AI custom training pipeline to perform full parameter fine-tuning on the policy documents daily."),
            ("B", "Build a Retrieval-Augmented Generation (RAG) pipeline by chunking the documents, generating embeddings, indexing them in Vertex AI Vector Search, and injecting retrieved contexts into the prompt."),
            ("C", "Train a custom BigQuery ML logistic regression model to predict mortgage text outputs."),
            ("D", "Train a student model using model distillation on Vertex AI.")
        ],
        "correct": "B",
        "explanation": "Retrieval-Augmented Generation (RAG) is the standard pattern to connect foundation models to dynamic, proprietary, and frequently updated external documents. It avoids the high computation cost of continuous model retraining and eliminates hallucinations by grounding the prompt with fresh retrieved context.",
        "ref": "Vertex AI - RAG Overview (https://cloud.google.com/vertex-ai/docs/generative-ai/rag-overview)"
    },
    {
        "category": "Practice Scenario 2 of 10",
        "title": "Hyperparameter Tuning for Diversity",
        "question": "You are building a creative marketing copy writer application powered by Gemini 1.5 Pro. The marketing team complains that the generated ad copy is too repetitive and predictable. How should you adjust the inference hyperparameters to increase the creativity and variety of the generated text?",
        "options": [
            ("A", "Lower the Temperature closer to 0.0 and set Top-K to 1."),
            ("B", "Increase the Temperature (e.g., to 1.0 or higher) and set a higher Top-P (e.g., 0.95) to allow a wider nucleus of candidate tokens."),
            ("C", "Increase the size of the model's context window."),
            ("D", "Freeze early transformer layers and perform LoRA fine-tuning.")
        ],
        "correct": "B",
        "explanation": "Temperature scales the next-token probability distribution; higher values increase randomness (creativity). Top-P (nucleus sampling) defines the cumulative probability boundary; raising it includes a wider set of candidate tokens, increasing diversity.",
        "ref": "Vertex AI - Configure Parameters (https://cloud.google.com/vertex-ai/docs/generative-ai/multimodal/configure-parameters)"
    },
    {
        "category": "Practice Scenario 3 of 10",
        "title": "Content Moderation & Safety Filters",
        "question": "A healthcare booking agent powered by the Gemini API must be deployed to the public. To ensure safety compliance, the agent must not generate dangerous medical advice or respond to sexually explicit inputs. How should you configure safety filters on Vertex AI?",
        "options": [
            ("A", "Build a custom BERT classifier container in GKE to scan all outputs."),
            ("B", "Configure the safetySettings block in the Gemini API request, mapping categories like HARM_CATEGORY_DANGEROUS_CONTENT to block thresholds (e.g., BLOCK_LOW_AND_ABOVE)."),
            ("C", "Wrap the model in a private VPC Service Controls perimeter."),
            ("D", "Perform supervised fine-tuning using toxic datasets.")
        ],
        "correct": "B",
        "explanation": "Vertex AI provides built-in safety filters for categories like Hate Speech, Harassment, Sexually Explicit, and Dangerous Content. Developers adjust block thresholds directly in the generation config API request to catch and filter violations.",
        "ref": "Vertex AI - Safety Settings (https://cloud.google.com/vertex-ai/docs/generative-ai/safety-settings)"
    },
    {
        "category": "Practice Scenario 4 of 10",
        "title": "Grounding with Enterprise Search",
        "question": "A company wants to build a search portal over thousands of internal product manuals without writing complex chunking, embedding, or database management code. Which Google Cloud service should they use?",
        "options": [
            ("A", "Vertex AI Custom Ingest pipelines on Dataproc."),
            ("B", "Vertex AI Agent Builder with a Data Store connected to their Cloud Storage bucket containing the manuals."),
            ("C", "BigQuery ML with external model references."),
            ("D", "Cloud Translation API connected to Pub/Sub.")
        ],
        "correct": "B",
        "explanation": "Vertex AI Agent Builder (specifically Enterprise Search) provides an out-of-the-box, no-code/low-code RAG system. By pointing a Data Store to a GCS bucket, the service automatically manages document ingestion, chunking, embeddings, and semantic grounding.",
        "ref": "Agent Builder - Create Data Store (https://cloud.google.com/generative-ai-app-builder/docs/create-datastore-search)"
    },
    {
        "category": "Practice Scenario 5 of 10",
        "title": "Parameter-Efficient Fine-Tuning",
        "question": "You need to adapt an open-source model (such as Gemma) to output proprietary medical code formats. You have limited GPU compute budget for training. Which training methodology is recommended on Vertex AI?",
        "options": [
            ("A", "Full parameter fine-tuning of all 7 billion model weights."),
            ("B", "Reinforcement Learning from Human Feedback (RLHF) from scratch."),
            ("C", "Parameter-Efficient Fine-Tuning (PEFT) using Low-Rank Adaptation (LoRA) to freeze base weights and train only small adapter matrices."),
            ("D", "Unsupervised pre-training on medical databases.")
        ],
        "correct": "C",
        "explanation": "PEFT/LoRA freezes the pre-trained model parameters and introduces small, trainable rank decomposition matrices. This dramatically reduces GPU memory requirements and training times, making custom domain training highly cost-efficient.",
        "ref": "Vertex AI - Tuning Overview (https://cloud.google.com/vertex-ai/docs/generative-ai/tuning/tuning-overview)"
    },
    {
        "category": "Practice Scenario 6 of 10",
        "title": "Model Distillation",
        "question": "An automotive company wants to run a natural language voice assistant inside cars. The hardware in the car has severe memory and processing limits and cannot host a large model. However, the assistant must mimic the reasoning quality of Gemini 1.5 Pro. How can Vertex AI solve this?",
        "options": [
            ("A", "Set up a high-performance TPU cluster inside the car dashboard."),
            ("B", "Execute a Model Distillation job on Vertex AI, using Gemini 1.5 Pro as the teacher to train a small, compact student model (like Gemma 2B) for edge deployment."),
            ("C", "Deploy a gRPC API connection requiring active satellite cellular networks."),
            ("D", "Convert the model weights to TF Lite format without retraining.")
        ],
        "correct": "B",
        "explanation": "Model distillation on Vertex AI uses a large, high-capability teacher model to train a smaller student model. The student model learns to approximate the teacher's outputs, enabling lightweight, low-latency deployments that retain high reasoning quality.",
        "ref": "Vertex AI - Model Distillation (https://cloud.google.com/vertex-ai/docs/generative-ai/models/distill-model)"
    },
    {
        "category": "Practice Scenario 7 of 10",
        "title": "Enforcing Structured JSON Outputs",
        "question": "A travel booking app uses Gemini to extract flight details (origin, destination, departure date) from customer emails. The extracted data must be fed into a legacy database that strictly accepts only valid JSON matching a specific schema. How should you invoke the Gemini API?",
        "options": [
            ("A", "Instruct the model in the prompt to write JSON and write a regex parser to fix errors."),
            ("B", "Use supervised fine-tuning to teach the model JSON syntax."),
            ("C", "Define a JSON schema using OpenAPI standards and pass it as the responseSchema configuration parameter in the generation config."),
            ("D", "Call Document AI Form Parser before prompting the model.")
        ],
        "correct": "C",
        "explanation": "The Gemini API supports schema enforcement. By passing an OpenAPI schema in responseSchema (and setting responseMimeType to application/json), the model is constrained to generate outputs that strictly comply with the target JSON structure, preventing parser errors.",
        "ref": "Vertex AI - Enforcing Schema (https://cloud.google.com/vertex-ai/docs/generative-ai/multimodal/control-structured-output)"
    },
    {
        "category": "Practice Scenario 8 of 10",
        "title": "IAM Permissions for Prompt Design",
        "question": "An external contractor joins your team to help prototype prompts and system instructions inside Vertex AI Studio. The contractor should not be able to train custom models, deploy endpoints, or access production billing. What IAM role is appropriate?",
        "options": [
            ("A", "roles/aiplatform.admin"),
            ("B", "roles/viewer"),
            ("C", "roles/aiplatform.user"),
            ("D", "roles/bigquery.dataEditor")
        ],
        "correct": "C",
        "explanation": "The roles/aiplatform.user role provides the necessary permissions to prototype prompts, call model prediction routes, and use Vertex AI Studio playgrounds, while restricting access to administrative settings.",
        "ref": "Vertex AI - Access Control (https://cloud.google.com/vertex-ai/docs/general/access-control)"
    },
    {
        "category": "Practice Scenario 9 of 10",
        "title": "Mitigating Summary Hallucinations",
        "question": "A hospital uses a model to summarize patient discharge notes. During trials, the model occasionally inserts fabricated medications not present in the patient's records. Which strategy is most effective to eliminate these hallucinations?",
        "options": [
            ("A", "Set the Temperature to 1.5 and increase Top-K."),
            ("B", "Apply supervised fine-tuning with a general vocabulary corpus."),
            ("C", "Ground the model by passing the complete medical record in the prompt context and setting the Temperature to 0.0 to enforce strict factual compliance."),
            ("D", "Switch to a smaller model size.")
        ],
        "correct": "C",
        "explanation": "Hallucinations are mitigated by grounding the model in the source context and setting the Temperature to 0.0 (or very low) to make token selection deterministic. This forces the model to synthesize text strictly from the provided inputs.",
        "ref": "Vertex AI - Grounding Generative Models (https://cloud.google.com/blog/products/ai-machine-learning/grounding-generative-ai-models-with-vertex-ai)"
    },
    {
        "category": "Practice Scenario 10 of 10",
        "title": "Optimizing Token Costs",
        "question": "You run a customer support chatbot that pre-loads a massive 50,000-token company FAQ context into every prompt. The chatbot processes thousands of queries hourly, leading to high API token costs. How can you optimize costs on Vertex AI?",
        "options": [
            ("A", "Shorten the FAQ list manually to 500 tokens."),
            ("B", "Enable Context Caching on Vertex AI to cache the FAQ token representations, paying a lower rate for cached tokens during subsequent queries."),
            ("C", "Migrate the model to run locally on a VM instance."),
            ("D", "Increase the Temperature to compress outputs.")
        ],
        "correct": "B",
        "explanation": "Context Caching allows caching of large, static context data (such as FAQs, system instructions, or document chunks) in the Gemini API. Subsequent calls referencing the cache pay a significantly lower rate for token processing, optimizing cost and reducing latency.",
        "ref": "Vertex AI - Context Caching Overview (https://cloud.google.com/vertex-ai/docs/generative-ai/context-caching-overview)"
    }
]

for s in scenarios:
    add_qa_slide(s["category"], s["title"], s["question"], s["options"], s["correct"], s["explanation"], s["ref"])

# Save PowerPoint deck
output_path = "Introduction_to_Generative_AI.pptx"
prs.save(output_path)
print(f"PowerPoint slide deck created successfully: {output_path}")
